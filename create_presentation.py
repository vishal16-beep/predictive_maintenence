from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLACK = RGBColor(0, 0, 0)
ORANGE = RGBColor(255, 107, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(20, 20, 20)
LIGHT_GRAY = RGBColor(200, 200, 200)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK

def add_title_box(slide, text, y=Inches(0.5)):
    txBox = slide.shapes.add_textbox(Inches(0.5), y, Inches(12), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ORANGE

def add_content_box(slide, text, y=Inches(1.8), size=Pt(24)):
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)

def add_orange_line(slide, y=Inches(1.4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(12), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Predictive Maintenance System", Inches(2))
add_content_box(slide, "NASA C-MAPSS FD001 Turbofan Degradation Analysis", Inches(3.2), Pt(32))
add_content_box(slide, "AI-Powered Remaining Useful Life Prediction", Inches(4.5), Pt(24))
add_orange_line(slide, Inches(3))

# Slide 2: Agenda
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Agenda")
add_orange_line(slide)
add_content_box(slide, """1. Problem Statement
2. Dataset Overview
3. Data Preprocessing
4. Feature Engineering
5. Model Development
6. API Development
7. Dashboard Development
8. Challenges & Solutions
9. Results & Demo""")

# Slide 3: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Problem Statement")
add_orange_line(slide)
add_content_box(slide, """Goal: Predict when a jet engine will fail before it happens

Why?
• Unplanned maintenance costs airlines $10-50B annually
• Safety risks from engine failures
• Inefficient maintenance scheduling

Solution: AI-powered Remaining Useful Life (RUL) prediction""")

# Slide 4: Dataset
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Dataset: NASA C-MAPSS FD001")
add_orange_line(slide)
add_content_box(slide, """Source: NASA Prognostics Center of Excellence

FD001 Dataset:
• 100 training engines, 100 test engines
• Degradation mode: High Pressure Compressor (HPC)
• Operating conditions: 6 flight conditions
• Sensors: 21 sensor readings per cycle

Files:
• train_FD001.txt - Training data (20,631 cycles)
• test_FD001.txt - Test data (13,096 cycles)
• RUL_FD001.txt - Actual RUL for test engines""")

# Slide 5: Data Preprocessing
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Data Preprocessing")
add_orange_line(slide)
add_content_box(slide, """Step 1: Data Loading
• Loaded space-separated files with 25 columns

Step 2: RUL Labeling
• RUL = max_cycle - current_cycle (clipped at 125)

Step 3: Drop Uninformative Sensors
• Removed: sensor_1, 5, 6, 10, 16, 18, 19
• Reason: Near-zero variance, no predictive power

Result: 17 informative features retained""")

# Slide 6: Feature Engineering
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Feature Engineering")
add_orange_line(slide)
add_content_box(slide, """17 Base Features:
• 3 Operational Settings
• 14 Informative Sensors

Rolling Window Features (window=10):
• Mean for each base feature (17)
• Standard deviation for each base feature (17)

Total: 51 engineered features

Train/Val Split:
• Split by engine ID (80/20)
• No row-level random split (prevents data leakage)""")

# Slide 7: Model Development
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Model Development")
add_orange_line(slide)
add_content_box(slide, """Algorithm: XGBoost Regressor

Hyperparameters:
• n_estimators: 500
• max_depth: 5
• learning_rate: 0.03
• subsample: 0.8
• colsample_bytree: 0.8
• min_child_weight: 5
• early_stopping_rounds: 30

Training: Early stopping on validation set""")

# Slide 8: NASA Scoring
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "NASA Asymmetric Scoring Function")
add_orange_line(slide)
add_content_box(slide, """Why asymmetric?
• Underestimating RUL is more dangerous than overestimating
• Late predictions cause failures
• Early predictions just waste maintenance budget

Formula:
• If underestimated (dangerous): exp(-diff/13) - 1
• If overestimated (safe): exp(diff/10) - 1

Penalizes late predictions 30% more heavily""")

# Slide 9: Model Performance
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Model Performance")
add_orange_line(slide)
add_content_box(slide, """Metric          Score           Interpretation
─────────────────────────────────────────────────
RMSE            19.44 cycles    Average error ~19 cycles
MAE             13.86 cycles    Mean absolute error
NASA Score      1338.5          Lower is better

Key Insights:
• Model captures degradation patterns well
• Early warnings possible with high confidence
• Suitable for production deployment""")

# Slide 10: API Development
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "API Development")
add_orange_line(slide)
add_content_box(slide, """Framework: FastAPI

Endpoints:
• GET  /health              → Health check
• GET  /api/v1/model/info   → Model metadata
• POST /api/v1/predict/rul  → Predict RUL

Request: JSON with sensor readings
Response: RUL prediction + urgency level

Features:
• Pydantic validation
• Async support
• Auto-generated docs""")

# Slide 11: Dashboard
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Dashboard Development")
add_orange_line(slide)
add_content_box(slide, """Framework: Streamlit + Plotly

Pages:
1. Predict RUL - Interactive sensor sliders
2. Batch Analysis - CSV upload
3. Model Info - Metrics & feature importance
4. Sensor Reference - Descriptions

Theme: NASA Aerospace (Black & Orange)
• Jet engine blueprint background
• High-visibility sliders
• Glowing orange accents""")

# Slide 12: Challenges
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Challenges & Solutions")
add_orange_line(slide)
add_content_box(slide, """Challenge                    Solution
─────────────────────────────────────────────────────
Data leakage                Split by engine ID, not rows
Uninformative sensors       Drop sensors with zero variance
Imbalanced RUL              Clip RUL at 125 cycles
Type mismatch in sliders    Convert all values to float
Background image loading    Use base64 encoding
Import path errors          Fixed config.settings imports""")

# Slide 13: Project Structure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Project Structure")
add_orange_line(slide)
add_content_box(slide, """predictive-maintenance/
├── config/                 # Settings
├── data/                   # NASA C-MAPSS files
├── src/
│   ├── data/              # Data loading
│   ├── inference/         # Prediction service
│   └── api/               # FastAPI app
├── models/                # Trained models
├── dashboard/             # Streamlit app
├── scripts/               # Training scripts
└── assets/                # Images""")

# Slide 14: Tech Stack
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Tech Stack")
add_orange_line(slide)
add_content_box(slide, """Component        Technology
─────────────────────────────────────
ML Model         XGBoost
Data Processing  Pandas, NumPy
API              FastAPI, Uvicorn
Dashboard        Streamlit, Plotly
Validation       Pydantic
Dataset          NASA C-MAPSS FD001""")

# Slide 15: Results
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Results Summary")
add_orange_line(slide)
add_content_box(slide, """✅ Completed:
• Real NASA dataset trained
• XGBoost model (RMSE: 19.44)
• FastAPI backend running
• Interactive dashboard
• Batch analysis capability
• NASA-themed UI

📊 Performance:
• 100 engines trained
• 51 engineered features
• Real-time predictions
• < 100ms response time""")

# Slide 16: Future
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Future Enhancements")
add_orange_line(slide)
add_content_box(slide, """1. Multi-fault support
   → FD002, FD003, FD004 datasets

2. Deep Learning models
   → LSTM/Transformer architectures

3. Real-time streaming
   → WebSocket integration

4. Cloud deployment
   → AWS/GCP/Azure

5. Mobile app
   → React Native

6. Digital twin
   → Engine simulation""")

# Slide 17: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_box(slide, "Thank You", Inches(2.5))
add_orange_line(slide, Inches(3.3))
add_content_box(slide, """Project: NASA C-MAPSS Predictive Maintenance

Stack: XGBoost + FastAPI + Streamlit

Dataset: NASA Commercial Modular Aero-Propulsion System Simulation

Questions?""", Inches(3.8), Pt(28))

prs.save('C:/Users/VISHAL/predictive-maintenance/PRESENTATION.pptx')
print("Presentation saved: PRESENTATION.pptx")
